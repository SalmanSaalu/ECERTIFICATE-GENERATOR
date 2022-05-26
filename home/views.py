
from django.shortcuts import render,redirect
from pptx import Presentation
import pandas as pd
from django.views import View
from .models import beforefetch

import os
import random

from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
from home.models import statusdetails
from django.contrib import messages
from django.utils.crypto import get_random_string
from .convter import ppt2pdf
import requests
from django.core.mail import send_mail, EmailMessage

def first(request):
    return render(request,'index.html')

class issueCertificate(View):
    def get(self,request):
        return render(request,'issuecertificate.html')
    def post(self,request):
        unique_id = get_random_string(length=32)
        a=request.FILES['excel']
        forcheckexcel=str(a)
        if forcheckexcel.endswith('.csv'):
            df = pd.read_csv(a)

            total_rows=len(df.axes[0]) #===> Axes of 0 is for a row
            total_cols=len(df.axes[1]) #===> Axes of 0 is for a column

            p=df.values.tolist()
            

            forexcel=[]
            for row in df:
                text = row[0:]
                forexcel.append(text)
            

            bc=request.FILES['template']
            forchecktemplate=str(bc)
            if forchecktemplate.endswith('.ppt') or forchecktemplate.endswith('.pptx'):
                    de=unique_id
                    beforefetch.objects.create(excel_file_upload=a,template_file_upload=bc,uniquename=de)
                    fortemplate=[]
                    ar1=[]
                    s=0
                    
                    prs = Presentation(bc)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if not shape.has_text_frame:
                                continue
                            text_frame = shape.text_frame
                            for paragraph in text_frame.paragraphs:
                                for run in paragraph.runs:
                                    a=str(run.text)
                                    
                                    b=len(a)
                                    
                                    for i in range(0,b):
                                            if(a[i]=="<" or s==10101):
      

                                                ar1.append(a[i])
                                                
                                                s=10101
                                                if(a[i]=='>'):
                                                    string=" "
                                                    joinedlist=string.join(ar1)
                                                    if '<' in joinedlist[1:-2] or '>' in joinedlist[1:-2]:
                                                                    messages.error(request,"** The '<' & '>' is not used in correct format")
                                                                    return render(request,'issuecertificate.html')
                                                    fortemplate.append(joinedlist)
                                                    ar1=[]
                                                    s=0
                    return render(request,'certificateparttwo.html',{'forexcel':forexcel,'fortemplate':fortemplate,'de':de})
            else:
                messages.error(request,'** the file you upoaded is not valid with the requirements needed')
                return render(request,'issuecertificate.html')
        else:
            messages.error(request,'** the file you upoaded is not valid with the requirements needed')
            return render(request,'issuecertificate.html')
        
class Certificateparttwo(View):
    def get(self,request):
        return redirect('/')
    def post(self,request):


        aa=request.POST['event']
        usr=beforefetch.objects.get(uniquename=aa)
        usr_excel=usr.excel_file_upload
        usr_template=usr.template_file_upload
        df = pd.read_csv(usr_excel)
        checkrow=[]
        for row in df:
            checkrow.append(row[0:])
        total_rows=len(df.axes[0])
        total_cols=len(df.axes[1])
        pst=[]
        pst2=[]
        fortemplate=[]
        ar1=[]
        s=0
        string=" "
        prs = Presentation(usr_template)
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                # elif shape.has_text_frame:
                #     pass
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        a=str(run.text)
                        
                        b=len(a)
                        
                        for i in range(0,b):
                                if(a[i]=="<" or s==10101):
                                    string=string+a[i]
                                    s=10101
                                    if(a[i]=='>'):
                                        fortemplate.append(string)
                                        string=" "
                                        s=0
        k=len(fortemplate)
        for i in range(1,k+1):
            aa=str(i)
            pst.append(aa)

        
        for i in pst:
            bb=request.POST[i]
            pst2.append(bb)
        ii=0
        jj=-1
        
        p=df.values.tolist()
        for ro in range(0,total_rows):
            
            emailing="hello"   
        
            for jjj in range(ro,total_rows):
                    for kkk in range(0,total_cols):
                         lll=p[jjj][kkk]
                         
                         lll=str(lll)
                         if '@gmail.com' in lll:
                            
                             emailing=lll
                             break
                    break
            if emailing=='hello':
                emailing='salmanpp05@gmail.com'

            prs = Presentation(usr_template)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            a=str(run.text)
                            
                            b=len(a)
                            
                            for i in range(0,b):
                                    if(a[i]=="<"):
                                        jj=jj+1
                                        w=pst2[jj]
                                        if w in checkrow:
                                            h=checkrow.index(w)
                                        p=df.values.tolist()
                                        hh=p[ii][h]
                                        u=len(str(hh))
                                        hh=str(hh)
                                        hh=str(hh)
                                        cu=" "
                                        cur_text = run.text
                                        print(fortemplate[jj])
                                        new_text = cur_text.replace(str(fortemplate[jj]), str(cu + hh))
                                        run.text = new_text
                                        print(run.text) 
            ii=ii+1
            jj=-1
            working_dir = os.getcwd()
            ww=random.randint(1,1000)
            # ww=str(ww)
            nn=str(ww)+'.pptx'
            aab=os.path.join( working_dir  + '/media/after_fetch',nn)
            # nn=str(ww)+'.pdf'
            # bba=os.path.join( working_dir  + '/media/after_fetch',nn)
            prs.save(os.path.abspath(aab))
            input_file_path = os.path.abspath(aab)
            # print(aab)
            # output_file_path = os.path.abspath(bba)
            # powerpoint =Dispatch("Powerpoint.Application",pythoncom.CoInitialize())
            # powerpoint.Visible = 1
            # slides = powerpoint.Presentations.Open(input_file_path)
            # slides.SaveAs(output_file_path, 32)
            # slides.Close()  
            # a=ppt2pdf(aab,nn)
            n=ppt2pdf(aab,nn)
            url="https://docs.google.com/presentation/d/"+n+"/export/pdf"
            print(url)
            # url="https://docs.google.com/document/export?format=[pdf]&id=[ID]"
            r = requests.get(url, allow_redirects=True)
            open('certificate.pdf', 'wb').write(r.content)



            mail=EmailMessage(
                        "certificate",
                        "user",
                        'salmansaalu10@gmail.com',
                        [emailing],

                    )

            mail.attach_file('certificate.pdf')
            mail.send(fail_silently=False,)  
            if emailing=='salmanpp05@gmail.com':
                emailing='email error'
                g=0
            else:
                g=1
            
            vv=statusdetails.objects.create(email=emailing,statusnamecheck=request.user.id,send=g)
            vv.save()
            
            
            os.remove('certificate.pdf')
            # fromaddr = "salmansaalu10@gmail.com"
            # toaddr = emailing
            #msg = MIMEMultipart()
            # msg['From'] = fromaddr
            # msg['To'] = toaddr 
            # msg['Subject'] = "Subject of the Mail"
            # body = "Body_of_the_mail"
            # msg.attach(MIMEText(body, 'plain'))
            # filename = nn
            # # attachment = open(output_file_path, "rb")
            # pp = MIMEBase('application', 'octet-stream')
            # pp.set_payload((attachment).read())
            # encoders.encode_base64(pp)
            # pp.add_header('Content-Disposition', "attachment; filename= %s" % filename)
            # msg.attach(pp)
            # s = smtplib.SMTP('smtp.gmail.com', 587)
            # s.starttls()
            # s.login(fromaddr, "yjiloplrlivfdsqg")
            # text = msg.as_string()
            # s.sendmail(fromaddr, toaddr, text)
            # s.quit()  
        return redirect('/')    



def status(request):
    pass
    ee=request.user.id
    ff=statusdetails.objects.filter(statusnamecheck=ee).order_by('-dateadded')
    return render(request,'status.html',{'statusdetails':ff})             
                                    
                                        
                                        
